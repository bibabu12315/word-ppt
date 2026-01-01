# -*- coding: utf-8 -*-
"""
文件名称：utils/create_template.py
主要作用：创建基础 PPT 模板
实现功能：
1. 使用 python-pptx 从零创建一个基础的 PPT 模板文件
2. 定义母版布局和占位符
3. 用于在没有外部模板时提供默认模板
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_demo_template(output_path: str):
    """
    创建一个符合新需求的测试用 PPT 模板 (4页结构)。
    1. Cover
    2. TOC
    3. Content Template
    4. End
    """
    prs = Presentation()
    slide_layout = prs.slide_layouts[6] # Blank

    # --- 1. Cover (Slide 0) ---
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tb.name = "cover_title"
    tb.text = "Cover Title Placeholder"
    tb.text_frame.paragraphs[0].font.size = Pt(44)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Meta info
    meta_fields = ["cover_company", "cover_project", "cover_presenter", "cover_dept", "cover_date"]
    for i, field in enumerate(meta_fields):
        tb = slide.shapes.add_textbox(Inches(1), Inches(4 + i*0.5), Inches(8), Inches(0.5))
        tb.name = field
        tb.text = field
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- 2. TOC (Slide 1) ---
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(1))
    tb.text = "Table of Contents"
    tb.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Prototypes for TOC items
    # page1_title_num (01)
    num_proto = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(1), Inches(0.5))
    num_proto.name = "page1_title_num"
    num_proto.text = "01"
    num_proto.text_frame.paragraphs[0].font.size = Pt(24)
    num_proto.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0) # Red for visibility
    
    # page1_title (Chapter Title)
    title_proto = slide.shapes.add_textbox(Inches(2.2), Inches(2), Inches(6), Inches(0.5))
    title_proto.name = "page1_title"
    title_proto.text = "Chapter Title Prototype"
    title_proto.text_frame.paragraphs[0].font.size = Pt(24)

    # --- 3. Content Template (Slide 2) ---
    slide = prs.slides.add_slide(slide_layout)
    
    # Nav Bar Prototype (page1_title)
    # Horizontal distribution
    nav_proto = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
    nav_proto.name = "page1_title"
    nav_proto.text = "Nav Item"
    nav_proto.text_frame.paragraphs[0].font.size = Pt(14)
    nav_proto.text_frame.paragraphs[0].font.bold = True
    
    # Description (page1_desc)
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    desc_box.name = "page1_desc"
    desc_box.text = "Description text goes here..."
    desc_box.text_frame.paragraphs[0].font.size = Pt(12)
    desc_box.text_frame.paragraphs[0].font.italic = True
    
    # Content Body (page1_bullet1)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(4))
    content_box.name = "page1_bullet1"
    content_box.text = "Content Body Placeholder"
    content_box.text_frame.paragraphs[0].font.size = Pt(18)

    # --- 4. End (Slide 3) ---
    slide = prs.slides.add_slide(slide_layout)
    tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tb.text = "Thank You"
    tb.text_frame.paragraphs[0].font.size = Pt(50)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Presenter on End Page
    tb_presenter = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    tb_presenter.name = "cover_presenter"
    tb_presenter.text = "Presenter Name"
    tb_presenter.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Save
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"New 4-slide template created at: {output_path}")
