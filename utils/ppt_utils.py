"""
文件名称：utils/ppt_utils.py
主要作用：PPT 操作工具函数库
实现功能：
1. 提供底层的 python-pptx 操作辅助函数
2. 实现幻灯片复制 (duplicate_slide)
3. 实现形状复制 (duplicate_shape)
4. 处理文本样式复制、图片复制等细节
"""
import copy
import uuid
from pptx.shapes.autoshape import Shape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.shapes.group import GroupShape
from pptx.shapes.connector import Connector
from pptx.oxml.ns import qn
from io import BytesIO

def duplicate_slide(pres, index):
    """
    Duplicate the slide at the given index in the presentation.
    Adds the new slide to the end of the presentation.
    Returns the new slide.
    """
    source_slide = pres.slides[index]
    # Use the same layout as the source slide to preserve background/master
    try:
        layout = source_slide.slide_layout
    except:
        # Fallback if layout access fails (rare)
        layout = pres.slide_layouts[0]
        
    dest_slide = pres.slides.add_slide(layout)

    # Map dest placeholders by idx
    dest_placeholders = {}
    for shape in dest_slide.placeholders:
        dest_placeholders[shape.placeholder_format.idx] = shape
        
    # List of shapes to remove from dest_slide (initially all, we will remove from this list if we keep them)
    shapes_to_remove = list(dest_slide.shapes)
    
    # Copy all shapes from source to dest
    for shape in source_slide.shapes:
        # Check if it's a placeholder
        if shape.is_placeholder:
            ph_format = shape.placeholder_format
            idx = ph_format.idx
            
            if idx in dest_placeholders:
                # Found matching placeholder!
                dest_ph = dest_placeholders[idx]
                
                # Don't remove it
                if dest_ph in shapes_to_remove:
                    shapes_to_remove.remove(dest_ph)
                
                # Copy content from source shape to dest placeholder
                _copy_placeholder_content(shape, dest_ph)
                
                # Ensure name matches (important for our logic)
                dest_ph.name = shape.name
                continue
        
        # If not a placeholder, or no match found, duplicate as new shape
        duplicate_shape(shape, dest_slide, keep_name=True)
        
    # Remove unused placeholders
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    
    return dest_slide

def _copy_placeholder_content(src, dst):
    """
    Copy content and properties from source placeholder to destination placeholder.
    Preserves the destination's inheritance link while copying overrides.
    """
    # Copy Text Body (preserves text and explicit formatting)
    if src.has_text_frame and dst.has_text_frame:
        if src.element.txBody is not None:
            new_txBody = copy.deepcopy(src.element.txBody)
            if dst.element.txBody is not None:
                idx = dst.element.index(dst.element.txBody)
                dst.element.remove(dst.element.txBody)
                dst.element.insert(idx, new_txBody)
            else:
                dst.element.append(new_txBody)
                
    # Copy Shape Properties (Fill, Line, etc.)
    if src.element.spPr is not None:
        new_spPr = copy.deepcopy(src.element.spPr)
        if dst.element.spPr is not None:
            idx = dst.element.index(dst.element.spPr)
            dst.element.remove(dst.element.spPr)
            dst.element.insert(idx, new_spPr)
        else:
            dst.element.insert(1, new_spPr)

    # Copy Shape Style (Theme references) for Placeholders
    # Use qn('p:style') to find the style element as it's not exposed as an attribute on CT_Shape
    src_style = src.element.find(qn('p:style'))
    dst_style = dst.element.find(qn('p:style'))
    
    if src_style is not None:
        new_style = copy.deepcopy(src_style)
        if dst_style is not None:
            idx = dst.element.index(dst_style)
            dst.element.remove(dst_style)
            dst.element.insert(idx, new_style)
        else:
            # Insert after spPr if possible, otherwise append
            if dst.element.spPr is not None:
                idx = dst.element.index(dst.element.spPr) + 1
                dst.element.insert(idx, new_style)
            else:
                dst.element.append(new_style)
    else:
        # If source has no style, remove destination style
        if dst_style is not None:
            dst.element.remove(dst_style)
            
    # Copy Geometry/Position overrides
    dst.left = src.left
    dst.top = src.top
    dst.width = src.width
    dst.height = src.height

def duplicate_shape(shape, slide, keep_name=False):
    """
    Duplicate a shape onto a specific slide.
    Preserves style (fill, line, effects) by copying spPr.
    """
    new_shape = None
    
    # 1. Text Box / AutoShape
    if isinstance(shape, Shape):
        try:
            autoshape_type_id = shape.auto_shape_type
            new_shape = slide.shapes.add_shape(
                autoshape_type_id,
                shape.left, shape.top, shape.width, shape.height
            )
        except ValueError:
            # Likely a TextBox
            new_shape = slide.shapes.add_textbox(
                shape.left, shape.top, shape.width, shape.height
            )

        # Copy Shape Properties (Fill, Line, Gradient, etc.)
        new_element = new_shape.element
        new_spPr = copy.deepcopy(shape.element.spPr)
        
        if new_element.spPr is not None:
            idx = new_element.index(new_element.spPr)
            new_element.remove(new_element.spPr)
            new_element.insert(idx, new_spPr)
        else:
            new_element.insert(1, new_spPr)

        # Copy Shape Style (Theme references)
        # Ensure we copy the style element if it exists, or remove it if it doesn't.
        # This fixes the issue where generated shapes have unwanted borders (from default theme style).
        src_style = shape.element.find(qn('p:style'))
        dst_style = new_element.find(qn('p:style'))
        
        if src_style is not None:
            new_style = copy.deepcopy(src_style)
            if dst_style is not None:
                idx = new_element.index(dst_style)
                new_element.remove(dst_style)
                new_element.insert(idx, new_style)
            else:
                # Insert after spPr
                # We know spPr exists because we just touched it.
                idx = new_element.index(new_element.spPr) + 1
                new_element.insert(idx, new_style)
        else:
            # Source has no style, remove destination style if present
            if dst_style is not None:
                new_element.remove(dst_style)

        # Copy text content
        if shape.has_text_frame and shape.element.txBody is not None:
            new_txBody = copy.deepcopy(shape.element.txBody)
            if new_txBody is not None:
                if new_element.txBody is not None:
                    idx = new_element.index(new_element.txBody)
                    new_element.remove(new_element.txBody)
                    new_element.insert(idx, new_txBody)
                else:
                    new_element.append(new_txBody)
        
    # 2. Picture
    elif isinstance(shape, Picture):
        try:
            image_stream = BytesIO(shape.image.blob)
            new_shape = slide.shapes.add_picture(
                image_stream, shape.left, shape.top, shape.width, shape.height
            )
            
            # Copy Picture Properties (e.g. cropping, effects)
            new_element = new_shape.element
            new_spPr = copy.deepcopy(shape.element.spPr)
            if new_element.spPr is not None:
                idx = new_element.index(new_element.spPr)
                new_element.remove(new_element.spPr)
                new_element.insert(idx, new_spPr)
                
        except Exception as e:
            print(f"Warning: Failed to copy picture {shape.name}: {e}")

    # 3. Connector
    elif isinstance(shape, Connector):
        try:
            new_shape = slide.shapes.add_connector(
                shape.connector_type, shape.begin_x, shape.begin_y, shape.end_x, shape.end_y
            )
            # Copy style
            new_element = new_shape.element
            new_spPr = copy.deepcopy(shape.element.spPr)
            if new_element.spPr is not None:
                idx = new_element.index(new_element.spPr)
                new_element.remove(new_element.spPr)
                new_element.insert(idx, new_spPr)
            else:
                new_element.insert(1, new_spPr)
        except Exception as e:
            print(f"Warning: Failed to copy connector {shape.name}: {e}")

    # 4. GroupShape or GraphicFrame (Fallback to XML cloning)
    else:
        try:
            # Clone the element directly
            new_element = copy.deepcopy(shape.element)
            
            # Ensure unique IDs
            for child in new_element.iter():
                if 'id' in child.attrib:
                     # Simple randomization to avoid collision
                     import random
                     child.set('id', str(random.randint(10000, 99999999)))
            
            slide.shapes._spTree.append(new_element)
            return None 
            
        except Exception as e:
            print(f"Warning: Failed to clone complex shape {shape.name}: {e}")

    # 5. Set Name (for supported types)
    if new_shape:
        if keep_name:
            new_shape.name = shape.name
        else:
            new_shape.name = f"{shape.name}_copy_{uuid.uuid4().hex[:8]}"
        
    return new_shape

def _copy_font(src_font, dest_font):
    """Helper to copy font attributes (Legacy, mostly unused if we copy txBody)"""
    if src_font.name:
        dest_font.name = src_font.name
    if src_font.size:
        dest_font.size = src_font.size
    if src_font.bold is not None:
        dest_font.bold = src_font.bold
    if src_font.italic is not None:
        dest_font.italic = src_font.italic
    if src_font.color and src_font.color.type == 1: # RGB
        dest_font.color.rgb = src_font.color.rgb

def move_slide(pres, old_index, new_index):
    """
    Move a slide from old_index to new_index.
    """
    xml_slides = pres.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])
