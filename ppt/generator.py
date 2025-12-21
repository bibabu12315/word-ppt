# -*- coding: utf-8 -*-
import os
import copy
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from parser.data_structs import PresentationData
from utils.ppt_utils import duplicate_slide, move_slide, duplicate_shape

class PPTGenerator:
    """
    PPT 生成器
    职责：读取 PPT 模板，根据 PresentationData 填充内容，保存为新文件。
    """

    def __init__(self, template_path: str, output_path: str):
        self.template_path = template_path
        self.output_path = output_path
        
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file not found: {template_path}")

    def generate(self, data: PresentationData):
        """
        执行生成过程
        """
        prs = Presentation(self.template_path)
        
        # 验证模板结构：必须至少有 4 页 (封面, 目录, 内容, 结尾)
        if len(prs.slides) < 4:
            print("Warning: Template should have at least 4 slides (Cover, TOC, Content, End).")
        
        # 1. 填充封面 (Slide 0)
        print("Generating Cover...")
        self._fill_cover(prs.slides[0], data)
        
        # 2. 填充目录 (Slide 1)
        print("Generating TOC...")
        self._fill_toc(prs.slides[1], data)
        
        # 3. 生成并填充内容页
        print("Generating Content Slides...")
        
        # 检查模板结构
        # 假设结构: 0:Cover, 1:TOC, 2:ChapterCover, 3:Content, 4:End
        has_chapter_cover = len(prs.slides) >= 5
        
        if not has_chapter_cover:
            print("Warning: Template does not seem to have a dedicated Chapter Cover slide (Index 2).")
            print("Assuming standard 4-slide structure. Chapter Cover will be skipped.")
            # Fallback to old logic if needed, or just error out?
            # For now, let's assume the user provided the correct template as requested.
            # But to be safe, we can just use Slide 2 as content if only 4 slides exist.
            tpl_cover_idx = -1
            tpl_content_idx = 2
            end_idx = 3
        else:
            tpl_cover_idx = 2
            tpl_content_idx = 3
            end_idx = 4
        
        num_chapters = len(data.slides)
        slide_pairs = [] # List of (cover_slide, content_slide)
        
        if num_chapters > 0:
            # 3.1 准备幻灯片对象
            if has_chapter_cover:
                # Chapter 1 uses the original templates
                slide_pairs.append((prs.slides[tpl_cover_idx], prs.slides[tpl_content_idx]))
                
                # Chapter 2+ duplicate the templates
                for i in range(1, num_chapters):
                    new_cover = duplicate_slide(prs, tpl_cover_idx)
                    new_content = duplicate_slide(prs, tpl_content_idx)
                    slide_pairs.append((new_cover, new_content))
                
                # Move End Slide to the very end
                move_slide(prs, end_idx, len(prs.slides) - 1)
                
            else:
                # Old logic: only content slides
                slide_pairs.append((None, prs.slides[tpl_content_idx]))
                for i in range(1, num_chapters):
                    new_content = duplicate_slide(prs, tpl_content_idx)
                    slide_pairs.append((None, new_content))
                move_slide(prs, end_idx, len(prs.slides) - 1)

            # 3.2 填充内容
            all_titles = [s.title for s in data.slides]
            
            # Page Index Counter
            # TOC is Page 1.
            # So next page (Chapter 1 Cover) is Page 2.
            current_page_idx = 2
            
            for i, chapter_data in enumerate(data.slides):
                cover_slide, content_slide = slide_pairs[i]
                
                # Fill Chapter Cover (if exists)
                if cover_slide:
                    self._fill_chapter_cover(cover_slide, chapter_data, i + 1)
                    # Chapter Cover counts as a page (Page 2), but doesn't show number
                    current_page_idx += 1
                
                # Fill Content Page
                if content_slide:
                    # Content Page is next (Page 3)
                    # Update chapter_data page_index for internal consistency if needed
                    chapter_data.page_index = current_page_idx
                    self._fill_content_page(content_slide, chapter_data, all_titles, i)
                    current_page_idx += 1
        
        # 3.3 填充结尾页 (End Slide)
        # End slide is now at the very end
        if len(prs.slides) > 0:
            self._fill_end_page(prs.slides[-1], data)
                
        # 4. 保存
        os.makedirs(os.path.dirname(self.output_path), exist_ok=True)
        prs.save(self.output_path)
        print(f"PPT generated successfully: {self.output_path}")

    def _find_shape(self, shape_map, base_name):
        """
        Helper to find a shape by base name, trying prefixes page1_, page2_, page3_.
        """
        prefixes = ["page1_", "page2_", "page3_"]
        
        # Try exact match first (if base_name already has prefix)
        if base_name in shape_map:
            return shape_map[base_name]
            
        # Try prefixes
        for prefix in prefixes:
            candidate = prefix + base_name
            if candidate in shape_map:
                return shape_map[candidate]
                
        return None

    def _fill_chapter_cover(self, slide, chapter_data, chapter_num):
        """
        填充章节封面页
        有且只有: page1_title (章节标题), page1_title_num (章节序号)
        无导航栏, 无页码
        """
        shape_map = self._build_shape_map(slide)
        
        # 1. 填充标题
        title_shape = self._find_shape(shape_map, "title")
        if title_shape:
            # 这里 page1_title 用作章节标题
            clean_title = self._clean_title(chapter_data.title)
            self._set_text(title_shape, clean_title)
            
        # 2. 填充序号
        num_shape = self._find_shape(shape_map, "title_num")
        if num_shape:
            num_text = f"{chapter_num:02d}"
            self._set_text(num_shape, num_text)
            
        # 3. 清理可能存在的其他干扰元素 (可选，如果模板很干净则不需要)
        # 用户说 "有且只有文本框...", 假设模板已经只保留了这两个。

    def _build_shape_map(self, slide) -> dict:
        """
        建立单个幻灯片的 shape_name -> shape 映射
        """
        mapping = {}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            mapping[shape.name] = shape
        return mapping

    def _estimate_text_width(self, text, font_size_emu):
        """
        估算文本宽度 (EMU)
        """
        if not text:
            return 0
        
        width_emu = 0
        for char in text:
            # 简单估算：汉字 1em，非汉字 0.6em
            if '\u4e00' <= char <= '\u9fff':
                width_emu += font_size_emu
            else:
                width_emu += font_size_emu * 0.6
        return int(width_emu)

    def _fill_cover(self, slide, data: PresentationData):
        """
        填充封面信息
        """
        shape_map = self._build_shape_map(slide)
        cover_mapping = {
            "cover_title": data.cover_title,
            "cover_company": data.meta_info.get("公司名称", ""),
            "cover_project": data.meta_info.get("项目名称", ""),
            "cover_presenter": data.meta_info.get("汇报人", ""),
            "cover_dept": data.meta_info.get("部门", "") or data.meta_info.get("部门 / 团队", ""),
            "cover_date": data.meta_info.get("日期", "")
        }

        for shape_name, text_content in cover_mapping.items():
            if shape_name in shape_map and text_content:
                self._set_text(shape_map[shape_name], text_content)

        # --- 封面排版优化：横向等距分布 ---
        # 需求：cover_dept -> 直接连接符 7 -> cover_presenter -> 直接连接符 6 -> cover_date
        # 位置：左下角 (以 cover_dept 的位置为基准)
        
        ordered_names = ["cover_dept", "直接连接符 7", "cover_presenter", "直接连接符 6", "cover_date"]
        valid_shapes = []
        
        # 更加鲁棒的查找逻辑 (处理可能的空格)
        for name in ordered_names:
            shape = shape_map.get(name)
            if not shape:
                # 尝试查找去除空格后的名称
                for s_name, s in shape_map.items():
                    if s_name.strip() == name:
                        shape = s
                        break
            if shape:
                valid_shapes.append(shape)
        
        if len(valid_shapes) > 1:
            # 1. 确定基准位置
            base_top = valid_shapes[0].top
            current_left = valid_shapes[0].left
            
            # 2. 设置间距 (1-2个空格，约 15pt)
            spacing = Pt(15)
            
            for shape in valid_shapes:
                tf = shape.text_frame
                
                # 设置不换行 & 自适应宽度
                tf.word_wrap = False
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                
                # 垂直对齐
                shape.top = base_top
                
                # 水平定位
                shape.left = current_left
                
                # 估算宽度并更新 current_left
                # 由于 python-pptx 无法实时获取渲染后的宽度，我们需要估算
                font_size_emu = Pt(18) # 默认 18pt
                try:
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        fs = tf.paragraphs[0].runs[0].font.size
                        if fs:
                            font_size_emu = fs
                except:
                    pass
                
                text = tf.text
                estimated_width = self._estimate_text_width(text, font_size_emu)
                
                # 加上间距
                current_left += estimated_width + spacing

    def _fill_toc(self, slide, data: PresentationData):
        """
        填充目录页
        动态生成: page1_title_num (01), page1_title (标题)
        """
        shape_map = self._build_shape_map(slide)
        
        # 查找原型 Shape
        proto_num = self._find_shape(shape_map, "title_num")
        proto_title = self._find_shape(shape_map, "title")
        
        if not proto_num or not proto_title:
            print("Warning: TOC prototypes (page1_title_num, page1_title) not found.")
            return

        # --- Dynamic Layout Calculation ---
        # Access presentation object via slide -> part -> package -> presentation_part -> presentation
        # Or simpler: since we created 'prs' in generate(), we could pass it.
        # But here we only have 'slide'.
        # In python-pptx, slide.part.package.presentation_part.presentation gives the Presentation object.
        # However, let's try a safer way if possible.
        # Actually, we can just use a fixed height if we can't get it, but getting it is better.
        try:
            prs = slide.part.package.presentation_part.presentation
            page_height = prs.slide_height
        except AttributeError:
            # Fallback: Assume standard 16:9 (10 inches height? No, 7.5 inches usually)
            # 7.5 inches = 6858000 EMUs
            page_height = 6858000 

        
        # Determine layout constraints
        start_top_num = proto_num.top
        start_top_title = proto_title.top
        
        # Use the lower starting point as the reference for "start_y"
        start_y = max(start_top_num, start_top_title)
        item_height = max(proto_num.height, proto_title.height)
        
        # Define bottom margin (use same as top margin for symmetry, or at least 1 inch)
        # Assuming 96 dpi, 1 inch = 914400 EMUs. 
        # Let's use a safe bottom margin.
        bottom_margin = start_y 
        max_y = page_height - bottom_margin
        
        num_items = len(data.slides)
        
        # Default spacing
        step_y = item_height * 1.5
        
        if num_items > 1:
            # Calculate available vertical span for the *starts* of the items
            # The last item starts at `max_y - item_height`
            available_span = max_y - start_y - item_height
            
            # Calculate required spacing to fit exactly
            calculated_step = available_span / (num_items - 1)
            
            # Use the smaller of (calculated_step, default_spacing) to avoid spreading too much
            # But if calculated_step is very small (negative even), we must compress.
            # So we actually want:
            # If calculated_step < default_spacing: use calculated_step (compress to fit)
            # If calculated_step > default_spacing: use default_spacing (don't spread too much)
            
            # However, we shouldn't overlap too much. 
            # Minimum spacing = item_height (touching)
            
            step_y = min(calculated_step, item_height * 1.5)
            
            # Ensure we don't overlap if possible (unless page is too small)
            if step_y < item_height:
                # Warning: Items will overlap. 
                # We could enforce step_y = item_height, but then it overflows page.
                # User asked to "limit them in the page", so we respect page bounds even if it overlaps.
                pass

        # --- Fix for Issue 2: Clone FIRST, then fill ---
        # We collect all shapes (original + clones) first, WITHOUT modifying them yet.
        # This ensures all clones are based on the clean prototype.
        
        toc_items = [] # List of (num_shape, title_shape)
        
        for i in range(len(data.slides)):
            if i == 0:
                # Use the prototype itself for the first item
                toc_items.append((proto_num, proto_title))
            else:
                # Clone the prototype (which is still clean because we haven't modified it yet)
                new_num = duplicate_shape(proto_num, slide)
                new_title = duplicate_shape(proto_title, slide)
                
                # Position
                offset = i * step_y
                new_num.top = start_top_num + int(offset)
                new_title.top = start_top_title + int(offset)
                
                toc_items.append((new_num, new_title))
        
        # Now fill text for all items
        for i, (num_shape, title_shape) in enumerate(toc_items):
            num_text = f"{i+1:02d}"
            # Clean title for TOC (remove "一、", "1." etc)
            title_text = self._clean_title(data.slides[i].title)
            
            self._set_text(num_shape, num_text)
            self._set_text(title_shape, title_text)
            
            # 规范化命名：page{i+2}_title (目录项对应正文页的标题)
            # 目录本身是 Page 1
            # Chapter 1 Cover is Page 2
            # Chapter 1 Content is Page 3
            # Chapter 2 Cover is Page 4
            # Chapter 2 Content is Page 5
            # Formula: target_page_idx = 2 + i * 2
            target_page_idx = 2 + i * 2
            title_shape.name = f"page{target_page_idx}_title"
            num_shape.name = f"page{target_page_idx}_title_num"

        # 4. 填充页码 (page1)
        self._ensure_page_number(slide, 1)

    def _fill_content_page(self, slide, chapter_data, all_titles, current_index):
        """
        填充内容页
        包含: 导航栏, 描述, 正文
        """
        shape_map = self._build_shape_map(slide)
        page_idx = chapter_data.page_index # 应该是 2, 3, 4...
        
        # 1. 生成导航栏 (page1_title) -> page{page_idx}_title
        proto_nav = self._find_shape(shape_map, "title")
        if proto_nav:
            # 获取页面尺寸
            prs = slide.part.package.presentation_part.presentation
            slide_height = prs.slide_height
            
            # 强制定位到左下角
            # 左边距 30pt, 下边距 30pt
            fixed_left = Pt(30)
            fixed_bottom = Pt(30)
            
            # 计算起始 Top (假设所有导航项高度一致，使用原型高度)
            start_top = slide_height - proto_nav.height - fixed_bottom
            start_left = fixed_left
            
            margin_x = proto_nav.width * 1.1 
            
            # --- Fix for Issue 4: Clone FIRST, then fill ---
            nav_items = []
            for i in range(len(all_titles)):
                # 计算当前项的位置
                current_item_left = start_left + int(i * margin_x)
                
                if i == 0:
                    shape = proto_nav
                else:
                    shape = duplicate_shape(proto_nav, slide)
                
                # 强制设置位置 (覆盖模板原有位置)
                shape.left = current_item_left
                shape.top = start_top
                
                nav_items.append(shape)
                
                # 规范化命名：导航栏的每一项对应一个页面
                # 第 i 项对应 Chapter i 的 Cover Page (Page 2 + i*2)
                target_page_idx = 2 + i * 2
                nav_items[-1].name = f"page{target_page_idx}_title"
            
            # Fill text and color
            for i, shape in enumerate(nav_items):
                # Clean title for Nav Bar
                clean_title = self._clean_title(all_titles[i])
                self._set_text(shape, clean_title)
                
                # Remove color override to respect template font settings
                # if i != current_index:
                #     self._set_font_color(shape, RGBColor(192, 192, 192)) # Gray
                # else:
                #     self._set_font_color(shape, RGBColor(0, 0, 0)) # Black
        else:
            print("Warning: Nav prototype (page1_title) not found on content slide.")

        # 2. 填充描述 (page1_desc)
        desc_shape = self._find_shape(shape_map, "desc")
        if desc_shape:
            self._set_text(desc_shape, chapter_data.description)

        # 3. 填充页码 (page1 -> pageX)
        self._ensure_page_number(slide, page_idx)

        # 4. 填充正文
        # 模式 A: 标题+内容 分离模式 (page1_bullet1 + page1_content1)
        bullet_shape = self._find_shape(shape_map, "bullet1")
        content_shape = self._find_shape(shape_map, "content1")
        
        if bullet_shape and content_shape:
            self._fill_content_paired(slide, bullet_shape, content_shape, chapter_data.blocks, shape_map, page_idx)
        
        # 模式 B: 仅有内容框 (page1_content1) - 自动分段
        elif content_shape:
            self._fill_content_multibox(slide, content_shape, chapter_data.blocks, shape_map, page_idx)

    def _fill_content_paired(self, slide, proto_title, proto_content, blocks, shape_map, page_idx):
        """
        分离模式：每个 Block 生成一对 (标题框, 内容框)
        - 标题框使用 pageX_bulletY
        - 内容框使用 pageX_contentY
        - 关键词框使用 pageX_keywordY
        """
        if not blocks:
            proto_title.text_frame.clear()
            proto_content.text_frame.clear()
            return

        # 查找关键词原型
        proto_keyword = self._find_shape(shape_map, "keyword1")

        # 布局参数
        start_top = proto_title.top
        
        # 计算标题和内容之间的垂直间距
        gap_title_content = max(Pt(5), proto_content.top - (proto_title.top + proto_title.height))
        
        # 段落之间的间距
        gap_paragraph = Pt(10)
        
        # 块与块之间的间距
        gap_block = Pt(20)
        
        current_top = start_top
        
        title_counter = 0
        content_counter = 0
        keyword_counter = 0
        
        for block in blocks:
            # 1. 处理标题 (Subtitle)
            title_counter += 1
            if title_counter == 1:
                title_shape = proto_title
            else:
                title_shape = duplicate_shape(proto_title, slide)
            
            # 规范化命名：page{page_idx}_bullet{title_counter}
            title_shape.name = f"page{page_idx}_bullet{title_counter}"
            
            title_shape.top = current_top
            self._set_text(title_shape, block.subtitle)
            current_top += title_shape.height + gap_title_content
            
            # 2. 处理内容 (Bullets)
            for bullet_text in block.bullets:
                content_counter += 1
                if content_counter == 1:
                    content_shape = proto_content
                else:
                    content_shape = duplicate_shape(proto_content, slide)
                
                # 规范化命名：page{page_idx}_content{content_counter}
                content_shape.name = f"page{page_idx}_content{content_counter}"
                
                content_shape.top = current_top
                self._set_text(content_shape, bullet_text)
                current_top += content_shape.height + gap_paragraph

            # 3. 处理关键词 (Keyword)
            if block.keyword and proto_keyword:
                keyword_counter += 1
                if keyword_counter == 1:
                    keyword_shape = proto_keyword
                else:
                    keyword_shape = duplicate_shape(proto_keyword, slide)
                
                # 规范化命名：page{page_idx}_keyword{keyword_counter}
                keyword_shape.name = f"page{page_idx}_keyword{keyword_counter}"
                
                keyword_shape.top = current_top
                self._set_text(keyword_shape, block.keyword)
                current_top += keyword_shape.height + gap_paragraph

            current_top += gap_block

    def _fill_content_multibox(self, slide, proto_shape, blocks, shape_map, page_idx):
        """
        新模式：将内容块拆分为多个文本框 (pageX_content1, pageX_content2...)
        """
        # 1. 收集所有需要显示的文本段落
        segments = []
        for block in blocks:
            # 标题作为单独一段
            if block.subtitle:
                segments.append({"text": block.subtitle, "is_title": True})
            # 列表/段落内容
            for bullet in block.bullets:
                segments.append({"text": bullet, "is_title": False})
        
        if not segments:
            proto_shape.text_frame.clear()
            return

        # 2. 布局参数
        start_top = proto_shape.top
        gap = Pt(10)
        current_top = start_top
        
        # 3. 生成文本框
        for i, seg in enumerate(segments):
            if i == 0:
                shape = proto_shape
            else:
                shape = duplicate_shape(proto_shape, slide)
            
            # 规范化命名：page{page_idx}_content{i+1}
            shape.name = f"page{page_idx}_content{i+1}"
            
            shape.top = current_top
            
            # 设置文本 (使用 _set_text 保留格式)
            self._set_text(shape, seg["text"])
            
            # 确保文本框高度自适应，以便正确计算堆叠位置
            shape.text_frame.word_wrap = True
            shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            # 更新位置
            current_top += shape.height + gap

    def _ensure_page_number(self, slide, page_idx):
        """
        确保页面右下角有页码
        """
        shape_map = self._build_shape_map(slide)
        target_name = f"page{page_idx}"
        
        # 1. 尝试查找现有的页码框 (可能是 page1 原型，或者是已经重命名的 pageX)
        page_shape = shape_map.get(target_name)
        if not page_shape:
            # 尝试查找原型 page1, page2, page3
            if "page1" in shape_map: page_shape = shape_map["page1"]
            elif "page2" in shape_map: page_shape = shape_map["page2"]
            elif "page3" in shape_map: page_shape = shape_map["page3"]
        
        # 2. 如果还没找到，创建一个新的
        if not page_shape:
            # 获取页面尺寸
            prs = slide.part.package.presentation_part.presentation
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # 尺寸和位置
            width = Pt(50)
            height = Pt(30)
            margin = Pt(20)
            
            left = slide_width - width - margin
            top = slide_height - height - margin
            
            page_shape = slide.shapes.add_textbox(left, top, width, height)
        
        # 3. 更新属性
        page_shape.name = target_name
        self._set_text(page_shape, str(page_idx))
        
        # 4. 强制定位到右下角 (用户要求：所有的页码文本框位置均在右下角)
        prs = slide.part.package.presentation_part.presentation
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # 保持原有宽高，或者给个最小值
        width = max(page_shape.width, Pt(30))
        height = max(page_shape.height, Pt(20))
        margin_right = Pt(30) # 右边距
        margin_bottom = Pt(20) # 下边距
        
        page_shape.left = slide_width - width - margin_right
        page_shape.top = slide_height - height - margin_bottom
        
        # 设置右对齐
        if page_shape.text_frame.paragraphs:
            page_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def _set_segment_text(self, shape, segment):
        """
        设置单个段落文本框的内容
        """
        text_frame = shape.text_frame
        text_frame.clear()
        
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = segment["text"]
        
        # 设置样式
        # 简单起见，标题加粗，正文普通
        # 实际应该参考模板样式，这里假设模板已经设置好了字体
        run.font.bold = segment["is_title"]
        
        # 尝试启用自动调整大小
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    def _set_blocks_text(self, shape, blocks):
        """
        将多个 ContentBlock 填充到一个文本框
        """
        text_frame = shape.text_frame
        
        # Capture style from first run of first paragraph
        template_run = None
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            template_run = text_frame.paragraphs[0].runs[0]
            
        text_frame.clear() # Clear all content
        
        for i, block in enumerate(blocks):
            # Subtitle
            if block.subtitle:
                p = text_frame.add_paragraph() if i > 0 or text_frame.paragraphs else text_frame.paragraphs[0]
                run = p.add_run()
                run.text = block.subtitle
                if template_run:
                    self._copy_font_style(template_run, run)
                run.font.bold = True
            
            # Bullets
            for bullet in block.bullets:
                p = text_frame.add_paragraph()
                p.level = 1
                run = p.add_run()
                run.text = bullet
                if template_run:
                    self._copy_font_style(template_run, run)

    def _set_text(self, shape, text):
        """
        设置文本，保留原有格式。
        Fix: 确保清除多余的段落，防止旧文本残留。
        Fix: 智能捕获样式 (从第一个非空 Run 继承)
        """
        text_frame = shape.text_frame
        
        if not text_frame.paragraphs:
            text_frame.add_paragraph()
            
        p = text_frame.paragraphs[0]
        
        # 1. Find a representative run to steal style from
        style_run = None
        if p.runs:
            # Use the first run that has text, or just the first run
            for r in p.runs:
                if r.text and r.text.strip():
                    style_run = r
                    break
            if not style_run and p.runs:
                style_run = p.runs[0]
        
        # 2. Prepare target run
        target_run = None
        if not p.runs:
            target_run = p.add_run()
        else:
            target_run = p.runs[0]
            
        # 3. Apply style from style_run to target_run (if they are different)
        if style_run and style_run != target_run:
            self._copy_font_style(style_run, target_run)
            
        # 4. Set text
        target_run.text = text
        
        # 5. Clear other runs in the first paragraph
        for i in range(1, len(p.runs)):
            p.runs[i].text = ""

        # 6. Remove extra paragraphs
        # If the original shape had multiple paragraphs, remove them.
        # We iterate backwards to avoid index issues.
        for i in range(len(text_frame.paragraphs) - 1, 0, -1):
            p_element = text_frame.paragraphs[i]._p
            p_element.getparent().remove(p_element)

    def _set_font_color(self, shape, rgb_color):
        """
        强制设置文本框内所有文字的颜色
        """
        if not shape.has_text_frame:
            return
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = rgb_color

    def _copy_font_style(self, src_run, dest_run):
        """
        将 src_run 的字体样式复制到 dest_run
        """
        if not src_run or not dest_run:
            return
        
        try:
            if src_run.font.name:
                dest_run.font.name = src_run.font.name
                self._set_ea_font(dest_run, src_run.font.name)
            
            if src_run.font.size:
                dest_run.font.size = src_run.font.size
            
            if src_run.font.bold is not None:
                dest_run.font.bold = src_run.font.bold
            if src_run.font.italic is not None:
                dest_run.font.italic = src_run.font.italic
                
            if src_run.font.color:
                if src_run.font.color.type == 1: # RGB
                    dest_run.font.color.rgb = src_run.font.color.rgb
                elif src_run.font.color.type == 2: # Theme Color
                    dest_run.font.color.theme_color = src_run.font.color.theme_color
                
        except Exception as e:
            print(f"Warning: Failed to copy font style: {e}")

    def _set_ea_font(self, run, font_name):
        """
        设置东亚字体
        """
        from pptx.oxml.ns import qn
        try:
            # Try standard access
            rPr = run._element.get_or_add_rPr()
        except AttributeError:
            # Fallback for _Run objects where _element might be named _r
            if hasattr(run, '_r'):
                rPr = run._r.get_or_add_rPr()
            else:
                return

        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'))
            rPr.append(ea)
        ea.set('typeface', font_name)

    def _fill_end_page(self, slide, data: PresentationData):
        """
        填充结尾页
        """
        shape_map = self._build_shape_map(slide)
        presenter = data.meta_info.get("汇报人", "")
        
        if "cover_presenter" in shape_map and presenter:
            self._set_text(shape_map["cover_presenter"], presenter)

    def _clean_title(self, title):
        """
        Remove leading numbering like "一、", "1.", "1、"
        """
        if not title:
            return ""
        # Pattern: Start of string, followed by (digits or chinese numerals), followed by dot or comma or chinese comma, optional whitespace
        pattern = r"^([0-9]+|[一二三四五六七八九十百]+)[、.．]\s*"
        return re.sub(pattern, "", title)
